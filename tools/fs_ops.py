import os
import shutil
import subprocess
from pathlib import Path
from typing import Any, Dict, List, Tuple

DEFAULT_IGNORE_DIRS = {
    ".git",
    ".venv",
    "venv",
    "node_modules",
    "__pycache__",
    ".mypy_cache",
    ".pytest_cache",
}


def resolve_safe_path(workspace_root: Path, raw_path: str) -> Path:
    if not isinstance(raw_path, str) or not raw_path.strip():
        raise ValueError("path must be a non-empty string")
    candidate = (workspace_root / raw_path).resolve()
    try:
        candidate.relative_to(workspace_root)
    except ValueError as err:
        raise ValueError(f"path escapes workspace: {raw_path}") from err
    return candidate


# Dangerous system paths that should NEVER be touched
_DANGEROUS_PATHS = {
    "C:\\Windows\\System32",
    "C:\\Windows\\SysWOW64",
    "/etc",
    "/sys",
    "/proc",
    "/boot",
    "/dev",
}


def resolve_any_path(raw_path: str) -> Path:
    """
    Resolve a path with FULL PC ACCESS - no workspace restriction.
    
    Used by FileAgent for unrestricted file operations across the entire PC.
    
    Rules:
    - Absolute paths → use directly
    - Paths starting with ~ → expand to home directory
    - Relative paths → resolve against user's home directory (NOT workspace)
    - Environment variables (%DESKTOP%, %USERPROFILE%, etc.) → expand automatically
    - SAFETY: Block genuinely dangerous system paths (Windows\\System32, /etc, /sys, /proc)
    
    Args:
        raw_path: The path string to resolve
        
    Returns:
        Resolved Path object
        
    Raises:
        ValueError: If path is empty or targets a dangerous system directory
    """
    if not isinstance(raw_path, str) or not raw_path.strip():
        raise ValueError("path must be a non-empty string")
    
    # Expand environment variables first (%DESKTOP%, %USERPROFILE%, etc.)
    expanded = os.path.expandvars(raw_path)
    
    # Convert to Path object
    candidate = Path(expanded)
    
    # If absolute, use it directly
    if candidate.is_absolute():
        resolved = candidate.resolve()
    # If starts with ~, expand to home
    elif expanded.startswith("~"):
        resolved = Path(expanded).expanduser().resolve()
    # Otherwise, resolve against home directory (not workspace)
    else:
        resolved = (Path.home() / expanded).resolve()
    
    # SAFETY CHECK: Block dangerous system paths
    resolved_str = str(resolved)
    for dangerous in _DANGEROUS_PATHS:
        if resolved_str.startswith(dangerous):
            raise ValueError(
                f"Access denied: {raw_path} targets protected system directory {dangerous}. "
                "FileAgent cannot modify OS internals."
            )
    
    return resolved


def to_rel(workspace_root: Path, path: Path) -> str:
    return str(path.relative_to(workspace_root)).replace("\\", "/")


def list_files(workspace_root: Path, path: str = ".", max_entries: int = 200, unrestricted: bool = False) -> Dict[str, Any]:
    """
    List files/directories.
    
    Args:
        workspace_root: Workspace root for safe path resolution
        path: Path to list (relative or absolute)
        max_entries: Maximum number of entries to return
        unrestricted: If True, allows access outside workspace (for FileAgent)
    """
    limit = max(1, min(int(max_entries), 1000))
    
    if unrestricted:
        root = resolve_any_path(path)
    else:
        root = resolve_safe_path(workspace_root, path)
    
    if not root.exists():
        raise FileNotFoundError(f"path not found: {path}")

    entries: List[Dict[str, Any]] = []
    if root.is_file():
        st = root.stat()
        return {
            "entries": [{"path": str(root) if unrestricted else to_rel(workspace_root, root), "type": "file", "size": st.st_size}],
            "truncated": False,
        }

    for current, dirnames, filenames in os.walk(root):
        dirnames[:] = [d for d in dirnames if d not in DEFAULT_IGNORE_DIRS]
        current_path = Path(current)
        for d in sorted(dirnames):
            p = current_path / d
            entries.append({"path": str(p) if unrestricted else to_rel(workspace_root, p), "type": "dir"})
            if len(entries) >= limit:
                return {"entries": entries, "truncated": True}
        for f in sorted(filenames):
            p = current_path / f
            try:
                size = p.stat().st_size
            except OSError:
                size = None
            entries.append({"path": str(p) if unrestricted else to_rel(workspace_root, p), "type": "file", "size": size})
            if len(entries) >= limit:
                return {"entries": entries, "truncated": True}
    return {"entries": entries, "truncated": False}


_MAX_READ_BYTES = 50 * 1024 * 1024  # 50 MB hard limit


def _detect_encoding(raw: bytes) -> str:
    """Try to detect the best text encoding for a byte sequence.
    Falls back through UTF-8 → cp1252 → latin-1 → replace."""
    # Check for null bytes — likely binary
    if b"\x00" in raw[:512]:
        raise ValueError("file appears to be binary (null bytes in first 512 bytes)")
    for enc in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            raw.decode(enc)
            return enc
        except (UnicodeDecodeError, LookupError):
            continue
    return "latin-1"  # guaranteed to decode any byte sequence


def read_file(workspace_root: Path, path: str, unrestricted: bool = False) -> Dict[str, Any]:
    """
    Read a text file.
    
    Args:
        workspace_root: Workspace root for safe path resolution
        path: Path to file (relative or absolute)
        unrestricted: If True, allows access outside workspace (for FileAgent)
    """
    if unrestricted:
        target = resolve_any_path(path)
    else:
        target = resolve_safe_path(workspace_root, path)
    
    if not target.exists() or not target.is_file():
        raise FileNotFoundError(f"file not found: {path}")

    # Guard: refuse files over 50 MB
    size = target.stat().st_size
    if size > _MAX_READ_BYTES:
        raise ValueError(
            f"file too large to read: {size / 1_048_576:.1f} MB "
            f"(limit is {_MAX_READ_BYTES // 1_048_576} MB). "
            "Use read_file_lines to read specific line ranges instead."
        )

    # Binary detection + encoding cascade
    raw = target.read_bytes()
    encoding = _detect_encoding(raw)
    text = raw.decode(encoding, errors="replace")
    return {"path": str(target) if unrestricted else to_rel(workspace_root, target), "content": text, "encoding": encoding}


def read_rich_file(workspace_root: Path, path: str, unrestricted: bool = False) -> Dict[str, Any]:
    """
    Read and extract text from rich file formats (PDF, DOCX, XLSX, images, etc.).
    
    Supported formats:
    - PDF: Extracts text using PyMuPDF (fitz)
    - DOCX: Extracts text using python-docx
    - XLSX/CSV: Returns formatted table using openpyxl/csv
    - PPTX: Extracts slide text using python-pptx
    - Images (JPG/PNG): Describes using vision API
    - ZIP: Lists contents
    - Plain text: Falls back to read_file
    
    Args:
        workspace_root: Workspace root for safe path resolution
        path: Path to file (relative or absolute)
        unrestricted: If True, allows access outside workspace (for FileAgent)
    
    Returns:
        Dict with 'path', 'content', 'format', and optional 'error' or 'install_hint'
    """
    if unrestricted:
        target = resolve_any_path(path)
    else:
        target = resolve_safe_path(workspace_root, path)
    
    if not target.exists() or not target.is_file():
        raise FileNotFoundError(f"file not found: {path}")
    
    ext = target.suffix.lower()
    result = {
        "path": str(target) if unrestricted else to_rel(workspace_root, target),
        "format": ext[1:] if ext else "unknown"
    }
    
    # PDF extraction
    if ext == ".pdf":
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(str(target))
            text_parts = []
            for page_num, page in enumerate(doc, 1):
                text = page.get_text()
                if text.strip():
                    text_parts.append(f"--- Page {page_num} ---\n{text}")
            doc.close()
            content = "\n\n".join(text_parts)
            result["content"] = content[:10000] + ("\n\n[truncated: content too long]" if len(content) > 10000 else "")
            result["pages"] = len(text_parts)
            return result
        except ImportError:
            result["error"] = "PyMuPDF not installed"
            result["install_hint"] = "To read PDF files, install: pip install PyMuPDF"
            return result
        except Exception as e:
            result["error"] = f"Failed to read PDF: {str(e)}"
            return result
    
    # DOCX extraction
    if ext == ".docx":
        try:
            from docx import Document
            doc = Document(str(target))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            content = "\n\n".join(paragraphs)
            result["content"] = content[:10000] + ("\n\n[truncated: content too long]" if len(content) > 10000 else "")
            result["paragraphs"] = len(paragraphs)
            return result
        except ImportError:
            result["error"] = "python-docx not installed"
            result["install_hint"] = "To read DOCX files, install: pip install python-docx"
            return result
        except Exception as e:
            result["error"] = f"Failed to read DOCX: {str(e)}"
            return result
    
    # XLSX extraction
    if ext in (".xlsx", ".xls"):
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(target), read_only=True, data_only=True)
            sheets_data = []
            for sheet_name in wb.sheetnames[:5]:  # Max 5 sheets
                sheet = wb[sheet_name]
                rows = []
                for row in list(sheet.iter_rows(values_only=True))[:100]:  # Max 100 rows per sheet
                    row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_str.strip():
                        rows.append(row_str)
                if rows:
                    sheets_data.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))
            wb.close()
            content = "\n\n".join(sheets_data)
            result["content"] = content[:10000] + ("\n\n[truncated: content too long]" if len(content) > 10000 else "")
            result["sheets"] = len(sheets_data)
            return result
        except ImportError:
            result["error"] = "openpyxl not installed"
            result["install_hint"] = "To read XLSX files, install: pip install openpyxl"
            return result
        except Exception as e:
            result["error"] = f"Failed to read XLSX: {str(e)}"
            return result
    
    # CSV extraction
    if ext == ".csv":
        try:
            import csv
            with open(target, 'r', encoding='utf-8', errors='replace') as f:
                reader = csv.reader(f)
                rows = []
                for i, row in enumerate(reader):
                    if i >= 100:  # Max 100 rows
                        break
                    row_str = " | ".join(row)
                    if row_str.strip():
                        rows.append(row_str)
            content = "\n".join(rows)
            result["content"] = content[:10000] + ("\n\n[truncated: content too long]" if len(content) > 10000 else "")
            result["rows"] = len(rows)
            return result
        except Exception as e:
            result["error"] = f"Failed to read CSV: {str(e)}"
            return result
    
    # PPTX extraction
    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(target))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
                if text_parts:
                    slides_text.append(f"--- Slide {i} ---\n" + "\n".join(text_parts))
            content = "\n\n".join(slides_text)
            result["content"] = content[:10000] + ("\n\n[truncated: content too long]" if len(content) > 10000 else "")
            result["slides"] = len(slides_text)
            return result
        except ImportError:
            result["error"] = "python-pptx not installed"
            result["install_hint"] = "To read PPTX files, install: pip install python-pptx"
            return result
        except Exception as e:
            result["error"] = f"Failed to read PPTX: {str(e)}"
            return result
    
    # Image description (requires vision API - placeholder for now)
    if ext in (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp"):
        result["content"] = f"[Image file: {target.name}]\nTo describe images, use the vision tools (capture_screen, read_screen_context) or open the file."
        result["note"] = "Image description requires vision API integration"
        return result
    
    # ZIP contents listing
    if ext == ".zip":
        try:
            import zipfile
            with zipfile.ZipFile(str(target), 'r') as zf:
                files = zf.namelist()[:100]  # Max 100 files
                content = "ZIP Contents:\n" + "\n".join(f"  - {f}" for f in files)
                if len(zf.namelist()) > 100:
                    content += f"\n\n[... and {len(zf.namelist()) - 100} more files]"
                result["content"] = content
                result["files"] = len(zf.namelist())
                return result
        except Exception as e:
            result["error"] = f"Failed to read ZIP: {str(e)}"
            return result
    
    # Plain text fallback - use regular read_file
    if ext in (".txt", ".md", ".py", ".js", ".ts", ".json", ".xml", ".html", ".css", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".log", ".sh", ".bat", ".ps1"):
        return read_file(workspace_root, path, unrestricted=unrestricted)
    
    # Unknown format
    result["error"] = f"Unsupported file format: {ext}"
    result["note"] = "Supported formats: PDF, DOCX, XLSX, CSV, PPTX, ZIP, and plain text files"
    return result


def search_text(workspace_root: Path, query: str, path: str = ".", max_results: int = 100, unrestricted: bool = False) -> Dict[str, Any]:
    """
    Search for text in files.
    
    Args:
        workspace_root: Workspace root for safe path resolution
        query: Text to search for
        path: Directory to search in
        max_results: Maximum number of results
        unrestricted: If True, allows access outside workspace (for FileAgent)
    """
    q = str(query).strip()
    if not q:
        raise ValueError("query is required")
    
    if unrestricted:
        root = resolve_any_path(path)
    else:
        root = resolve_safe_path(workspace_root, path)
    
    if not root.exists():
        raise FileNotFoundError(f"path not found: {path}")

    limit = max(1, min(int(max_results), 300))

    rg = shutil.which("rg")
    if rg:
        cmd = [rg, "-n", "--no-heading", "--color", "never", "--hidden"]
        for ignore_dir in sorted(DEFAULT_IGNORE_DIRS):
            cmd.extend(["--glob", f"!{ignore_dir}/**"])
        cmd.extend([q, str(root)])
        out = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
        if out.returncode not in (0, 1):
            raise RuntimeError(out.stderr.strip() or "rg failed")
        lines = [ln for ln in out.stdout.splitlines() if ln.strip()][:limit]
        return {"matches": lines, "truncated": len(lines) >= limit, "engine": "rg"}

    matches: List[str] = []
    if root.is_file():
        files = [root]
    else:
        files = []
        for p in root.rglob("*"):
            if not p.is_file():
                continue
            if any(part in DEFAULT_IGNORE_DIRS for part in p.parts):
                continue
            files.append(p)
    for file_path in files:
        try:
            for idx, line in enumerate(file_path.read_text(encoding="utf-8", errors="ignore").splitlines(), 1):
                if q in line:
                    path_str = str(file_path) if unrestricted else to_rel(workspace_root, file_path)
                    matches.append(f"{path_str}:{idx}:{line}")
                    if len(matches) >= limit:
                        return {"matches": matches, "truncated": True, "engine": "python"}
        except OSError:
            continue
    return {"matches": matches, "truncated": False, "engine": "python"}


def _audit_write(path: str, status: str, byte_count: int) -> None:
    """Append a write event to ~/.ankita/audit.log for forensics (never crashes)."""
    try:
        import time as _time
        audit_dir = Path.home() / ".ankita"
        audit_dir.mkdir(parents=True, exist_ok=True)
        log_line = (
            f"[{_time.strftime('%Y-%m-%d %H:%M:%S')}] "
            f"Attempted write to: {path} | Status: {status} | Bytes: {byte_count}\n"
        )
        with open(audit_dir / "audit.log", "a", encoding="utf-8") as f:
            f.write(log_line)
    except Exception:
        pass


def write_file(workspace_root: Path, path: str, content: str, overwrite: bool = True) -> Dict[str, Any]:
    """
    Write content to a file.

    Supports ABSOLUTE paths (e.g. C:\\Users\\anime\\Desktop\\poem.txt) directly —
    these bypass workspace restriction so agents can save to Desktop.
    For relative paths, resolves against workspace_root as before.

    Returns a RECEIPT dict — agent must see status:'success' before claiming the file was saved.
    """
    raw_path = str(path).strip()
    candidate = Path(raw_path)

    # Absolute path (Desktop saves, etc.) — use directly without workspace restriction
    if candidate.is_absolute():
        target = candidate.resolve()
    else:
        target = resolve_safe_path(workspace_root, raw_path)

    absolute_path = str(target)
    _audit_write(absolute_path, "attempting", 0)

    existed = target.exists()
    if existed and not bool(overwrite):
        _audit_write(absolute_path, "FAILED: file exists and overwrite=false", 0)
        raise FileExistsError(f"file exists and overwrite=false: {path}")

    target.parent.mkdir(parents=True, exist_ok=True)
    byte_count = len(str(content).encode("utf-8"))
    try:
        target.write_text(str(content), encoding="utf-8")
        _audit_write(absolute_path, "success", byte_count)
    except Exception as err:
        _audit_write(absolute_path, f"FAILED: {err}", 0)
        raise

    return {
        "status": "success",
        "path": absolute_path,
        "absolute_path": absolute_path,
        "FILE_PATH": absolute_path,   # picked up by _extract_artifacts
        "bytes": byte_count,
        "overwrote": existed,
    }


def edit_file(
    workspace_root: Path,
    path: str,
    old_text: str,
    new_text: str,
    replace_all: bool = False,
) -> Dict[str, Any]:
    if old_text == "":
        raise ValueError("old_text must be non-empty")
    target = resolve_safe_path(workspace_root, path)
    if not target.exists() or not target.is_file():
        raise FileNotFoundError(f"file not found: {path}")

    text = target.read_text(encoding="utf-8")
    if old_text not in text:
        raise ValueError("old_text not found in file")

    if replace_all:
        updated = text.replace(old_text, new_text)
        count = text.count(old_text)
    else:
        updated = text.replace(old_text, new_text, 1)
        count = 1
    target.write_text(updated, encoding="utf-8")
    return {"path": to_rel(workspace_root, target), "replacements": count}


def delete_path(workspace_root: Path, path: str, recursive: bool = False, missing_ok: bool = False, unrestricted: bool = False) -> Dict[str, Any]:
    """
    Delete a file or directory.
    
    Args:
        workspace_root: Workspace root for safe path resolution
        path: Path to delete
        recursive: Allow deleting directories
        missing_ok: Don't error if path doesn't exist
        unrestricted: If True, allows access outside workspace (for FileAgent)
    """
    if unrestricted:
        target = resolve_any_path(path)
    else:
        target = resolve_safe_path(workspace_root, path)
    
    if not target.exists():
        if missing_ok:
            return {"path": path, "deleted": False, "reason": "not_found"}
        raise FileNotFoundError(f"path not found: {path}")

    if target.is_dir():
        if not recursive:
            raise IsADirectoryError("target is a directory; set recursive=true")
        shutil.rmtree(target)
        return {"path": str(target) if unrestricted else to_rel(workspace_root, target), "deleted": True, "type": "dir"}

    target.unlink()
    return {"path": str(target) if unrestricted else to_rel(workspace_root, target), "deleted": True, "type": "file"}


def rename_path(workspace_root: Path, path: str, new_name: str, overwrite: bool = False) -> Dict[str, Any]:
    source = resolve_safe_path(workspace_root, path)
    if not source.exists():
        raise FileNotFoundError(f"path not found: {path}")
    n = str(new_name).strip()
    if not n or "/" in n or "\\" in n:
        raise ValueError("new_name must be a simple filename")

    dest = source.with_name(n)
    resolve_safe_path(workspace_root, str(dest.relative_to(workspace_root)))
    if dest.exists() and not overwrite:
        raise FileExistsError("destination exists; set overwrite=true")
    if dest.exists() and overwrite:
        if dest.is_dir():
            shutil.rmtree(dest)
        else:
            dest.unlink()
    source.rename(dest)
    return {"from": to_rel(workspace_root, source), "to": to_rel(workspace_root, dest)}


def move_path(workspace_root: Path, src: str, dst: str, overwrite: bool = False, unrestricted: bool = False) -> Dict[str, Any]:
    """
    Move/rename a file or directory.
    
    Args:
        workspace_root: Workspace root for safe path resolution
        src: Source path
        dst: Destination path
        overwrite: Allow overwriting existing destination
        unrestricted: If True, allows access outside workspace (for FileAgent)
    """
    if unrestricted:
        source = resolve_any_path(src)
        dest = resolve_any_path(dst)
    else:
        source = resolve_safe_path(workspace_root, src)
        dest = resolve_safe_path(workspace_root, dst)
    
    if not source.exists():
        raise FileNotFoundError(f"path not found: {src}")
    if dest.exists() and not overwrite:
        raise FileExistsError("destination exists; set overwrite=true")
    dest.parent.mkdir(parents=True, exist_ok=True)
    if dest.exists() and overwrite:
        if dest.is_dir():
            shutil.rmtree(dest)
        else:
            dest.unlink()
    shutil.move(str(source), str(dest))
    
    if unrestricted:
        return {"from": str(source), "to": str(dest)}
    else:
        return {"from": to_rel(workspace_root, source), "to": to_rel(workspace_root, dest)}


def copy_path(workspace_root: Path, src: str, dst: str, overwrite: bool = False, recursive: bool = False, unrestricted: bool = False) -> Dict[str, Any]:
    """
    Copy a file or directory.
    
    Args:
        workspace_root: Workspace root for safe path resolution
        src: Source path
        dst: Destination path
        overwrite: Allow overwriting existing destination
        recursive: Allow copying directories
        unrestricted: If True, allows access outside workspace (for FileAgent)
    """
    if unrestricted:
        source = resolve_any_path(src)
        dest = resolve_any_path(dst)
    else:
        source = resolve_safe_path(workspace_root, src)
        dest = resolve_safe_path(workspace_root, dst)
    
    if not source.exists():
        raise FileNotFoundError(f"path not found: {src}")
    if dest.exists() and not overwrite:
        raise FileExistsError("destination exists; set overwrite=true")
    dest.parent.mkdir(parents=True, exist_ok=True)

    if source.is_dir():
        if not recursive:
            raise IsADirectoryError("source is directory; set recursive=true")
        if dest.exists() and overwrite:
            shutil.rmtree(dest)
        shutil.copytree(source, dest)
        if unrestricted:
            return {"from": str(source), "to": str(dest), "type": "dir"}
        else:
            return {"from": to_rel(workspace_root, source), "to": to_rel(workspace_root, dest), "type": "dir"}

    if dest.exists() and overwrite:
        dest.unlink()
    shutil.copy2(source, dest)
    if unrestricted:
        return {"from": str(source), "to": str(dest), "type": "file"}
    else:
        return {"from": to_rel(workspace_root, source), "to": to_rel(workspace_root, dest), "type": "file"}


def make_dir(workspace_root: Path, path: str, parents: bool = True, exist_ok: bool = True, unrestricted: bool = False) -> Dict[str, Any]:
    """
    Create a directory.
    
    Args:
        workspace_root: Workspace root for safe path resolution
        path: Directory path to create
        parents: Create parent directories if needed
        exist_ok: Don't error if directory already exists
        unrestricted: If True, allows creation outside workspace (for FileAgent)
    """
    if unrestricted:
        target = resolve_any_path(path)
    else:
        target = resolve_safe_path(workspace_root, path)
    
    target.mkdir(parents=bool(parents), exist_ok=bool(exist_ok))
    
    if unrestricted:
        return {"path": str(target), "created": True, "absolute_path": str(target)}
    else:
        return {"path": to_rel(workspace_root, target), "created": True}


def file_info(workspace_root: Path, path: str, unrestricted: bool = False) -> Dict[str, Any]:
    """
    Get file/directory metadata.
    
    Args:
        workspace_root: Workspace root for safe path resolution
        path: Path to inspect
        unrestricted: If True, allows access outside workspace (for FileAgent)
    """
    if unrestricted:
        target = resolve_any_path(path)
    else:
        target = resolve_safe_path(workspace_root, path)
    
    if not target.exists():
        raise FileNotFoundError(f"path not found: {path}")
    st = target.stat()
    return {
        "path": str(target) if unrestricted else to_rel(workspace_root, target),
        "type": "dir" if target.is_dir() else "file",
        "size": st.st_size,
        "mtime": st.st_mtime,
    }


def read_file_lines(
    workspace_root: Path,
    path: str,
    start_line: int,
    end_line: int,
) -> Dict[str, Any]:
    """
    Read a specific range of lines from a file without loading the entire file.

    Lines are 1-indexed and inclusive on both ends, matching what editors show.
    Useful for the "Look Before You Leap" pattern — read the exact broken lines
    before calling edit_file_lines.

    Args:
        workspace_root: Workspace root for path resolution.
        path:           Path to the file (absolute or relative to workspace).
        start_line:     First line to return (1-indexed, inclusive).
        end_line:       Last line to return (1-indexed, inclusive).

    Returns:
        Dict with ok, path, start_line, end_line, content (str), line_count.
    """
    try:
        fpath = resolve_safe_path(workspace_root, path)
        if not fpath.is_file():
            return {"ok": False, "error": f"Not a file: {path}"}

        all_lines = fpath.read_text(encoding="utf-8", errors="replace").splitlines(keepends=True)
        total = len(all_lines)

        # Clamp to valid range (1-indexed → 0-indexed)
        s = max(1, start_line)
        e = min(end_line, total)
        if s > total:
            return {
                "ok": False,
                "error": f"start_line {start_line} exceeds file length {total}",
            }

        selected = all_lines[s - 1: e]
        # Prefix each line with its line number for easy LLM reference
        numbered = "".join(f"{s + i}: {line}" for i, line in enumerate(selected))

        return {
            "ok": True,
            "path": to_rel(workspace_root, fpath),
            "start_line": s,
            "end_line": e,
            "total_lines": total,
            "content": numbered,
            "line_count": len(selected),
        }
    except Exception as err:
        return {"ok": False, "error": str(err)}


def edit_file_lines(
    workspace_root: Path,
    path: str,
    start_line: int,
    end_line: int,
    new_content: str,
) -> Dict[str, Any]:
    """
    Surgically replace a specific range of lines in a file.

    Only the targeted lines are changed — all surrounding code is left intact.
    Lines are 1-indexed and inclusive on both ends.

    IMPORTANT: Always call read_file or read_file_lines first to confirm the
    exact line numbers before calling this function — never guess.

    Args:
        workspace_root: Workspace root for path resolution.
        path:           Path to the file (absolute or relative to workspace).
        start_line:     First line to replace (1-indexed, inclusive).
        end_line:       Last line to replace (1-indexed, inclusive).
        new_content:    Replacement text. May contain multiple lines separated
                        by \\n. A trailing newline will be added automatically
                        if missing so the file structure stays intact.

    Returns:
        Dict with ok, path, lines_replaced, total_lines_before, total_lines_after.
    """
    try:
        fpath = resolve_safe_path(workspace_root, path)
        if not fpath.is_file():
            return {"ok": False, "error": f"Not a file: {path}"}

        all_lines = fpath.read_text(encoding="utf-8", errors="replace").splitlines(keepends=True)
        total_before = len(all_lines)

        # Validate range
        s = max(1, start_line)
        e = min(end_line, total_before)
        if s > total_before:
            return {
                "ok": False,
                "error": f"start_line {start_line} exceeds file length {total_before}",
            }
        if s > e:
            return {
                "ok": False,
                "error": f"start_line {s} is greater than end_line {e}",
            }

        # Ensure new_content ends with a newline so surrounding lines stay intact
        if new_content and not new_content.endswith("\n"):
            new_content = new_content + "\n"

        # Split new_content into lines, preserving newlines
        replacement_lines = [ln if ln.endswith("\n") else ln + "\n"
                             for ln in new_content.splitlines()]

        # Perform the surgical replacement (0-indexed slice)
        all_lines[s - 1: e] = replacement_lines

        fpath.write_text("".join(all_lines), encoding="utf-8")
        total_after = len(all_lines)

        return {
            "ok": True,
            "path": to_rel(workspace_root, fpath),
            "start_line": s,
            "end_line": e,
            "lines_replaced": e - s + 1,
            "replacement_line_count": len(replacement_lines),
            "total_lines_before": total_before,
            "total_lines_after": total_after,
        }
    except Exception as err:
        return {"ok": False, "error": str(err)}


def check_syntax(workspace_root: Path, path: str) -> Dict[str, Any]:
    """
    Run a fast Python syntax and AST check on a file using the built-in
    compiler — no subprocess, no external linter required.

    Use this AFTER editing code with edit_file_lines, BEFORE running it
    with run_command. Catches indentation errors, SyntaxErrors, and
    malformed expressions instantly.

    Args:
        workspace_root: Workspace root for path resolution.
        path:           Path to the Python file to check.

    Returns:
        Dict with ok (bool), errors (list of error strings), and a summary.
    """
    try:
        import ast as _ast
        fpath = resolve_safe_path(workspace_root, path)

        if not fpath.is_file():
            return {"ok": False, "errors": [f"File not found: {path}"], "summary": "File not found."}

        if fpath.suffix.lower() not in {".py", ".pyw"}:
            return {
                "ok": False,
                "errors": [f"Not a Python file: {path}"],
                "summary": "check_syntax only supports .py/.pyw files.",
            }

        source = fpath.read_text(encoding="utf-8", errors="replace")
        errors = []

        # Step 1: compile() catches SyntaxError + IndentationError
        try:
            compile(source, str(fpath), "exec")
        except SyntaxError as e:
            errors.append(
                f"SyntaxError at line {e.lineno}: {e.msg} — {(e.text or '').strip()}"
            )

        # Step 2: ast.parse() catches additional structural issues
        if not errors:
            try:
                _ast.parse(source, filename=str(fpath))
            except SyntaxError as e:
                errors.append(
                    f"AST SyntaxError at line {e.lineno}: {e.msg} — {(e.text or '').strip()}"
                )

        if errors:
            return {
                "ok": False,
                "path": to_rel(workspace_root, fpath),
                "errors": errors,
                "summary": f"❌ {len(errors)} syntax error(s) found — fix before running.",
            }

        return {
            "ok": True,
            "path": to_rel(workspace_root, fpath),
            "errors": [],
            "summary": f"✅ Syntax OK — {fpath.name} is clean and ready to run.",
        }

    except Exception as err:
        return {"ok": False, "errors": [str(err)], "summary": f"check_syntax failed: {err}"}


def apply_patch(workspace_root: Path, patch: str) -> Dict[str, Any]:
    begin = "*** Begin Patch"
    end = "*** End Patch"
    add = "*** Add File: "
    delete = "*** Delete File: "
    update = "*** Update File: "
    move = "*** Move to: "

    raw = patch.strip()
    if not raw:
        raise ValueError("patch is empty")
    lines = raw.splitlines()
    if lines[0].strip() != begin or lines[-1].strip() != end:
        raise ValueError("patch must start with '*** Begin Patch' and end with '*** End Patch'")

    i = 1
    summary = {"added": [], "modified": [], "deleted": []}

    def is_hunk_header(line: str) -> bool:
        s = line.strip()
        return s.startswith(add) or s.startswith(delete) or s.startswith(update)

    while i < len(lines) - 1:
        line = lines[i]
        stripped = line.strip()
        if not stripped:
            i += 1
            continue

        if stripped.startswith(add):
            path = stripped[len(add) :].strip()
            target = resolve_safe_path(workspace_root, path)
            content_lines: List[str] = []
            i += 1
            while i < len(lines) - 1 and not is_hunk_header(lines[i].strip()):
                row = lines[i]
                if not row.startswith("+"):
                    raise ValueError(f"invalid add line for {path}: '{row}'")
                content_lines.append(row[1:])
                i += 1
            target.parent.mkdir(parents=True, exist_ok=True)
            target.write_text("\n".join(content_lines) + ("\n" if content_lines else ""), encoding="utf-8")
            summary["added"].append(to_rel(workspace_root, target))
            continue

        if stripped.startswith(delete):
            path = stripped[len(delete) :].strip()
            target = resolve_safe_path(workspace_root, path)
            if target.exists():
                if target.is_dir():
                    shutil.rmtree(target)
                else:
                    target.unlink()
            summary["deleted"].append(path.replace("\\", "/"))
            i += 1
            continue

        if stripped.startswith(update):
            path = stripped[len(update) :].strip()
            src = resolve_safe_path(workspace_root, path)
            if not src.exists() or not src.is_file():
                raise FileNotFoundError(f"file not found for update: {path}")

            i += 1
            move_to: str | None = None
            if i < len(lines) - 1 and lines[i].strip().startswith(move):
                move_to = lines[i].strip()[len(move) :].strip()
                i += 1

            op_lines: List[str] = []
            while i < len(lines) - 1 and not is_hunk_header(lines[i].strip()):
                row = lines[i]
                if row.strip() == "":
                    op_lines.append(" ")
                    i += 1
                    continue
                if row.startswith("@@"):
                    op_lines.append(row)
                elif row.startswith(" ") or row.startswith("+") or row.startswith("-"):
                    op_lines.append(row)
                else:
                    raise ValueError(f"invalid update line for {path}: '{row}'")
                i += 1

            updated = _apply_update_ops(src.read_text(encoding="utf-8"), op_lines, path)
            if move_to:
                dst = resolve_safe_path(workspace_root, move_to)
                dst.parent.mkdir(parents=True, exist_ok=True)
                dst.write_text(updated, encoding="utf-8")
                src.unlink()
                summary["modified"].append(to_rel(workspace_root, dst))
            else:
                src.write_text(updated, encoding="utf-8")
                summary["modified"].append(to_rel(workspace_root, src))
            continue

        raise ValueError(f"invalid patch hunk header: '{line}'")

    return summary


def _apply_update_ops(original_text: str, op_lines: List[str], path: str) -> str:
    lines = original_text.splitlines()
    trailing_newline = original_text.endswith("\n")

    cursor = 0
    idx = 0
    while idx < len(op_lines):
        if op_lines[idx].startswith("@@"):
            idx += 1
            continue

        block: List[str] = []
        while idx < len(op_lines) and not op_lines[idx].startswith("@@"):
            block.append(op_lines[idx])
            idx += 1

        old_lines: List[str] = []
        new_lines: List[str] = []
        for row in block:
            prefix = row[:1]
            body = row[1:] if row else ""
            if prefix == " ":
                old_lines.append(body)
                new_lines.append(body)
            elif prefix == "-":
                old_lines.append(body)
            elif prefix == "+":
                new_lines.append(body)
            else:
                raise ValueError(f"invalid update marker in {path}: '{row}'")

        if not old_lines and not new_lines:
            continue

        if not old_lines:
            lines[cursor:cursor] = new_lines
            cursor += len(new_lines)
            continue

        pos = _find_subsequence(lines, old_lines, cursor)
        if pos is None:
            raise ValueError(f"failed to find expected block in {path}: {old_lines[:3]}")
        lines[pos : pos + len(old_lines)] = new_lines
        cursor = pos + len(new_lines)

    out = "\n".join(lines)
    if trailing_newline:
        out += "\n"
    return out


def _find_subsequence(lines: List[str], pattern: List[str], start: int) -> int | None:
    if not pattern:
        return start
    limit = len(lines) - len(pattern)
    for i in range(max(start, 0), limit + 1):
        if lines[i : i + len(pattern)] == pattern:
            return i
    for i in range(max(start, 0), limit + 1):
        if [ln.rstrip() for ln in lines[i : i + len(pattern)]] == [p.rstrip() for p in pattern]:
            return i
    return None


# ─────────────────────────────────────────────────────────────────────────────
# PHASE 2 TOOLS — Advanced FileAgent capabilities
# ─────────────────────────────────────────────────────────────────────────────

def pc_search(query: str, file_types: List[str] = None, max_results: int = 50) -> Dict[str, Any]:
    """
    Search for files across the entire PC by name pattern.
    
    Args:
        query: Search pattern (supports wildcards like *.txt)
        file_types: Optional list of extensions to filter (e.g. ['.pdf', '.docx'])
        max_results: Maximum number of results to return
        
    Returns:
        Dict with list of matching file paths
    """
    import glob
    
    results = []
    search_locations = [
        Path.home() / "Desktop",
        Path.home() / "Documents",
        Path.home() / "Downloads",
        Path.home() / "Pictures",
    ]
    
    for location in search_locations:
        if not location.exists():
            continue
        
        try:
            # Search recursively
            pattern = f"**/*{query}*" if not query.startswith("*") else f"**/{query}"
            for path in location.glob(pattern):
                if len(results) >= max_results:
                    break
                
                # Filter by file type if specified
                if file_types and path.suffix.lower() not in [ft.lower() for ft in file_types]:
                    continue
                
                results.append({
                    "path": str(path),
                    "name": path.name,
                    "size": path.stat().st_size if path.is_file() else 0,
                    "modified": path.stat().st_mtime if path.exists() else 0,
                })
        except (PermissionError, OSError):
            continue
    
    return {
        "query": query,
        "count": len(results),
        "results": results[:max_results],
        "truncated": len(results) > max_results,
    }


def trash_path(workspace_root: Path, path: str, unrestricted: bool = False) -> Dict[str, Any]:
    """
    Move a file or directory to the recycle bin instead of permanent deletion.
    
    Args:
        workspace_root: Workspace root for safe path resolution
        path: Path to move to trash
        unrestricted: If True, allows trashing files outside workspace
        
    Returns:
        Dict with status
    """
    if unrestricted:
        target = resolve_any_path(path)
    else:
        target = resolve_safe_path(workspace_root, path)
    
    if not target.exists():
        raise FileNotFoundError(f"path not found: {path}")
    
    # Use send2trash library if available, otherwise use OS-specific commands
    try:
        import send2trash
        send2trash.send2trash(str(target))
        return {
            "path": str(target),
            "trashed": True,
            "method": "send2trash",
        }
    except ImportError:
        # Fallback to OS-specific commands
        if os.name == "nt":
            # Windows: use PowerShell to move to recycle bin
            import subprocess
            cmd = f'Add-Type -AssemblyName Microsoft.VisualBasic; [Microsoft.VisualBasic.FileIO.FileSystem]::DeleteFile("{target}", "OnlyErrorDialogs", "SendToRecycleBin")'
            result = subprocess.run(
                ["powershell", "-Command", cmd],
                capture_output=True,
                text=True,
            )
            if result.returncode == 0:
                return {
                    "path": str(target),
                    "trashed": True,
                    "method": "powershell",
                }
            else:
                raise RuntimeError(f"Failed to trash: {result.stderr}")
        else:
            # Unix: move to ~/.local/share/Trash
            trash_dir = Path.home() / ".local" / "share" / "Trash" / "files"
            trash_dir.mkdir(parents=True, exist_ok=True)
            dest = trash_dir / target.name
            shutil.move(str(target), str(dest))
            return {
                "path": str(target),
                "trashed": True,
                "method": "unix_trash",
                "trash_location": str(dest),
            }


def disk_analysis(workspace_root: Path, path: str = ".", unrestricted: bool = False) -> Dict[str, Any]:
    """
    Analyze disk usage by directory.
    
    Args:
        workspace_root: Workspace root for safe path resolution
        path: Directory to analyze
        unrestricted: If True, allows analyzing directories outside workspace
        
    Returns:
        Dict with disk usage breakdown
    """
    if unrestricted:
        target = resolve_any_path(path)
    else:
        target = resolve_safe_path(workspace_root, path)
    
    if not target.is_dir():
        raise ValueError(f"path is not a directory: {path}")
    
    def get_dir_size(dir_path: Path) -> int:
        """Calculate total size of directory."""
        total = 0
        try:
            for item in dir_path.rglob("*"):
                if item.is_file():
                    try:
                        total += item.stat().st_size
                    except (PermissionError, OSError):
                        pass
        except (PermissionError, OSError):
            pass
        return total
    
    # Analyze immediate subdirectories
    subdirs = []
    total_size = 0
    
    try:
        for item in target.iterdir():
            if item.is_dir():
                size = get_dir_size(item)
                subdirs.append({
                    "name": item.name,
                    "path": str(item),
                    "size_bytes": size,
                    "size_mb": round(size / (1024 * 1024), 2),
                })
                total_size += size
            elif item.is_file():
                size = item.stat().st_size
                total_size += size
    except (PermissionError, OSError) as e:
        return {"error": f"Permission denied: {e}"}
    
    # Sort by size descending
    subdirs.sort(key=lambda x: x["size_bytes"], reverse=True)
    
    return {
        "path": str(target),
        "total_size_bytes": total_size,
        "total_size_mb": round(total_size / (1024 * 1024), 2),
        "total_size_gb": round(total_size / (1024 * 1024 * 1024), 2),
        "subdirectories": subdirs[:20],  # Top 20 largest
        "count": len(subdirs),
    }


def diff_files(workspace_root: Path, file1: str, file2: str, unrestricted: bool = False) -> Dict[str, Any]:
    """
    Compare two text files and show differences.
    
    Args:
        workspace_root: Workspace root for safe path resolution
        file1: First file path
        file2: Second file path
        unrestricted: If True, allows comparing files outside workspace
        
    Returns:
        Dict with diff output
    """
    if unrestricted:
        path1 = resolve_any_path(file1)
        path2 = resolve_any_path(file2)
    else:
        path1 = resolve_safe_path(workspace_root, file1)
        path2 = resolve_safe_path(workspace_root, file2)
    
    if not path1.exists():
        raise FileNotFoundError(f"file1 not found: {file1}")
    if not path2.exists():
        raise FileNotFoundError(f"file2 not found: {file2}")
    
    # Read files
    try:
        with open(path1, "r", encoding="utf-8") as f:
            lines1 = f.readlines()
    except UnicodeDecodeError:
        return {"error": f"file1 is not a text file: {file1}"}
    
    try:
        with open(path2, "r", encoding="utf-8") as f:
            lines2 = f.readlines()
    except UnicodeDecodeError:
        return {"error": f"file2 is not a text file: {file2}"}
    
    # Generate unified diff
    import difflib
    diff = list(difflib.unified_diff(
        lines1,
        lines2,
        fromfile=str(path1),
        tofile=str(path2),
        lineterm="",
    ))
    
    # Count changes
    additions = sum(1 for line in diff if line.startswith("+") and not line.startswith("+++"))
    deletions = sum(1 for line in diff if line.startswith("-") and not line.startswith("---"))
    
    return {
        "file1": str(path1),
        "file2": str(path2),
        "identical": len(diff) == 0,
        "additions": additions,
        "deletions": deletions,
        "diff": "\n".join(diff[:500]),  # Limit to 500 lines
        "truncated": len(diff) > 500,
    }


def bulk_op(workspace_root: Path, operation: str, paths: List[str], destination: str = None, unrestricted: bool = False) -> Dict[str, Any]:
    """
    Perform batch operations on multiple files.
    
    Args:
        workspace_root: Workspace root for safe path resolution
        operation: Operation to perform (move, copy, delete, trash)
        paths: List of file paths to operate on
        destination: Destination directory (for move/copy operations)
        unrestricted: If True, allows operations outside workspace
        
    Returns:
        Dict with operation results
    """
    if operation not in ("move", "copy", "delete", "trash"):
        raise ValueError(f"Invalid operation: {operation}. Must be one of: move, copy, delete, trash")
    
    if operation in ("move", "copy") and not destination:
        raise ValueError(f"destination is required for {operation} operation")
    
    results = []
    succeeded = 0
    failed = 0
    
    for path_str in paths:
        try:
            if operation == "delete":
                result = delete_path(workspace_root, path_str, recursive=True, unrestricted=unrestricted)
                results.append({"path": path_str, "status": "deleted", "ok": True})
                succeeded += 1
            elif operation == "trash":
                result = trash_path(workspace_root, path_str, unrestricted=unrestricted)
                results.append({"path": path_str, "status": "trashed", "ok": True})
                succeeded += 1
            elif operation == "move":
                # Construct full destination path with filename
                filename = Path(path_str).name
                dest_path = str(Path(destination) / filename)
                result = move_path(workspace_root, path_str, dest_path, overwrite=True, unrestricted=unrestricted)
                results.append({"path": path_str, "status": "moved", "destination": dest_path, "ok": True})
                succeeded += 1
            elif operation == "copy":
                # Construct full destination path with filename
                filename = Path(path_str).name
                dest_path = str(Path(destination) / filename)
                result = copy_path(workspace_root, path_str, dest_path, overwrite=True, recursive=True, unrestricted=unrestricted)
                results.append({"path": path_str, "status": "copied", "destination": dest_path, "ok": True})
                succeeded += 1
        except Exception as e:
            results.append({"path": path_str, "status": "failed", "error": str(e), "ok": False})
            failed += 1
    
    return {
        "operation": operation,
        "total": len(paths),
        "succeeded": succeeded,
        "failed": failed,
        "results": results,
    }
