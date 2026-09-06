"""Parsers: pdf->md (pymupdf4llm), docx/pptx/xlsx/html/md/txt/eml; OCR fallback.

Every parser returns (markdown, n_pages, kind). Corrupt files raise —
callers catch and mark status=failed (never crash chat).
"""

from __future__ import annotations

import html as _html
import re
from html.parser import HTMLParser
from pathlib import Path

SUPPORTED = {".pdf", ".docx", ".pptx", ".xlsx", ".html", ".htm", ".md", ".markdown",
             ".txt", ".eml", ".csv"}


class _TextExtract(HTMLParser):
    SKIP = {"script", "style", "nav", "footer", "header", "aside", "form", "noscript"}

    def __init__(self):
        super().__init__()
        self._skip = 0
        self._chunks: list[str] = []

    def handle_starttag(self, tag, attrs):
        if tag.lower() in self.SKIP:
            self._skip += 1
        elif tag.lower() in ("p", "br", "h1", "h2", "h3", "h4", "li", "tr", "title"):
            self._chunks.append("\n")

    def handle_endtag(self, tag):
        if tag.lower() in self.SKIP and self._skip:
            self._skip -= 1
        else:
            self._chunks.append("\n")

    def handle_data(self, data):
        if not self._skip and data and data.strip():
            self._chunks.append(data)

    def text(self) -> str:
        blob = "".join(self._chunks)
        blob = _html.unescape(blob)
        blob = re.sub(r"[ \t]+", " ", blob)
        return re.sub(r"\n\s*\n+", "\n\n", blob).strip()


def html_to_text(html_text: str) -> str:
    try:
        p = _TextExtract()
        p.feed(html_text or "")
        return p.text()
    except Exception:
        return re.sub(r"\s+", " ", re.sub(r"<[^>]+>", " ", html_text or "")).strip()


def parse_file(path: str | Path) -> tuple[str, int, str]:
    p = Path(path)
    suf = p.suffix.lower()
    if suf == ".pdf":
        return _parse_pdf(p)
    if suf == ".docx":
        return _parse_docx(p)
    if suf == ".pptx":
        return _parse_pptx(p)
    if suf in (".xlsx", ".csv"):
        return _parse_sheet(p)
    if suf in (".html", ".htm"):
        return _parse_html(p)
    if suf == ".eml":
        return _parse_eml(p)
    return _parse_text(p)


def _parse_pdf(p: Path) -> tuple[str, int, str]:
    import os
    md, n_pages = "", 1
    try:
        import fitz
        doc = fitz.open(str(p))
        n_pages = max(1, doc.page_count)
        doc.close()
    except Exception:
        pass
    try:
        import pymupdf4llm
        md = pymupdf4llm.to_markdown(str(p)) or ""
    except Exception:
        md = ""
    if not md.strip() and os.getenv("ZUMBA_VAULT_OCR", "1") != "0":
        md = _ocr_pdf(p, n_pages) or md
    if not md.strip():
        raise ValueError(f"empty PDF extraction: {p.name}")
    return md.strip(), n_pages, "pdf"


def _ocr_pdf(p: Path, n_pages: int) -> str:
    try:
        import fitz
        doc = fitz.open(str(p))
        bits = []
        for i, page in enumerate(doc):
            try:
                tp = page.get_textpage_ocr(dpi=200) if hasattr(page, "get_textpage_ocr") else None
                t = (tp.extractText() if tp is not None else "") or ""
            except Exception:
                t = ""
            if t.strip():
                bits.append(f"\n\n## Page {i + 1} (OCR)\n\n{t.strip()}")
        doc.close()
        return "\n".join(bits)
    except Exception:
        return ""


def _parse_docx(p: Path) -> tuple[str, int, str]:
    from docx import Document
    doc = Document(str(p))
    bits = []
    for para in doc.paragraphs:
        t = (para.text or "").strip()
        if not t:
            continue
        style = str(getattr(getattr(para, "style", None), "name", "") or "")
        if style.startswith("Heading"):
            try:
                lvl = int("".join(c for c in style if c.isdigit()) or 1)
            except Exception:
                lvl = 1
            bits.append(f"\n\n{'#' * max(1, min(lvl, 4))} {t}\n")
        else:
            bits.append(t)
    for table in getattr(doc, "tables", []):
        rows = []
        for row in table.rows:
            rows.append("| " + " | ".join((c.text or "").strip() for c in row.cells) + " |")
        if rows:
            bits.append("\n\n" + "\n".join(rows) + "\n")
    return ("\n\n".join(bits).strip() or p.stem, 1, "docx")


def _parse_pptx(p: Path) -> tuple[str, int, str]:
    from pptx import Presentation
    prs = Presentation(str(p))
    bits = []
    for i, slide in enumerate(prs.slides):
        bits.append(f"\n\n## Slide {i + 1}\n")
        for shape in slide.shapes:
            t = (getattr(shape, "text", "") or "").strip()
            if t:
                bits.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    bits.append("| " + " | ".join((c.text or "").strip() for c in row.cells) + " |")
    return ("\n\n".join(bits).strip() or p.stem, max(1, len(prs.slides)), "pptx")


def _parse_sheet(p: Path) -> tuple[str, int, str]:
    if p.suffix.lower() == ".csv":
        text = p.read_text(encoding="utf-8", errors="replace")
        lines = [l for l in text.splitlines() if l.strip()][:200]
        return ("# " + p.stem + "\n\n" + "\n".join("| " + l.replace(",", " | ") + " |" for l in lines), 1, "csv")
    import openpyxl
    wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
    bits = []
    for ws in wb.worksheets:
        bits.append(f"\n\n## Sheet: {ws.title}\n")
        for row in ws.iter_rows(values_only=True):
            vals = [(str(v).strip() if v is not None else "") for v in row]
            if any(vals):
                bits.append("| " + " | ".join(vals) + " |")
    return ("\n".join(bits).strip() or p.stem, 1, "xlsx")


def _parse_html(p: Path) -> tuple[str, int, str]:
    text = html_to_text(p.read_text(encoding="utf-8", errors="replace"))
    return (f"# {p.stem}\n\n{text}".strip(), 1, "html")


def _parse_eml(p: Path) -> tuple[str, int, str]:
    import email
    from email import policy
    raw = p.read_bytes()
    msg = email.message_from_bytes(raw, policy=policy.default)
    subject = str(msg.get("Subject", "") or "")
    frm = str(msg.get("From", "") or "")
    date = str(msg.get("Date", "") or "")
    body = ""
    try:
        if msg.is_multipart():
            for part in msg.walk():
                ctype = part.get_content_type()
                if ctype == "text/plain":
                    body += str(part.get_content() or "") + "\n"
        else:
            body = str(msg.get_content() or "")
    except Exception:
        body = ""
    md = f"# {subject or p.stem}\n\nFrom: {frm}\nDate: {date}\n\n{(body or '').strip()}"
    return (md.strip(), 1, "eml")


def _parse_text(p: Path) -> tuple[str, int, str]:
    text = p.read_text(encoding="utf-8", errors="replace")
    if p.suffix.lower() in (".md", ".markdown"):
        return (text.strip() or p.stem, 1, "md")
    return ((f"# {p.stem}\n\n{text}").strip(), 1, "txt")
